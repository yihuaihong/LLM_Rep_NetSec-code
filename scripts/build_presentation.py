"""Build the final NeurIPS presentation deck (13 slides, 16:9, ~15 min talk).

Replaces the existing optimistic deck with a deck that reflects the actual
negative findings + V1-V4 validation protocol contribution.

Output: LLM_RepSpace_NetSec_Presentation.pptx
"""
import os, json
import numpy as np
import matplotlib
matplotlib.use("Agg")
import matplotlib.pyplot as plt
from pathlib import Path
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR

BASE = Path(__file__).resolve().parent.parent
METRICS = BASE / "results" / "metrics"
FIGS_DIR = BASE / "paper" / "figures"
DECK_FIGS = BASE / "paper" / "deck_figures"
DECK_FIGS.mkdir(parents=True, exist_ok=True)

# ------------ Color theme ------------
NAVY    = RGBColor(0x1F, 0x3A, 0x5F)   # title color
ACCENT  = RGBColor(0xE6, 0x4D, 0x4D)   # red accent (failure / negative)
GREEN   = RGBColor(0x2E, 0xA8, 0x5C)   # success accent (positive control)
GRAY_T  = RGBColor(0x4A, 0x4A, 0x4A)   # body text
LIGHT_BG = RGBColor(0xF5, 0xF7, 0xFA)  # light background panels
DARK_BG  = RGBColor(0xFF, 0xFF, 0xFF)


def add_textbox(slide, left, top, width, height, text, *,
                font_size=18, bold=False, color=GRAY_T, align=PP_ALIGN.LEFT,
                font_name="Calibri"):
    tb = slide.shapes.add_textbox(Inches(left), Inches(top), Inches(width), Inches(height))
    tf = tb.text_frame
    tf.word_wrap = True
    tf.margin_left = Emu(0)
    tf.margin_right = Emu(0)
    tf.margin_top = Emu(0)
    tf.margin_bottom = Emu(0)
    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.name = font_name
    run.font.size = Pt(font_size)
    run.font.bold = bold
    run.font.color.rgb = color
    return tb


def add_bullets(slide, left, top, width, height, items, *,
                font_size=16, color=GRAY_T, font_name="Calibri",
                bullet_char="•"):
    tb = slide.shapes.add_textbox(Inches(left), Inches(top), Inches(width), Inches(height))
    tf = tb.text_frame
    tf.word_wrap = True
    tf.margin_left = Emu(0); tf.margin_right = Emu(0)
    tf.margin_top = Emu(0);  tf.margin_bottom = Emu(0)
    for i, item in enumerate(items):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        p.space_before = Pt(2); p.space_after = Pt(4)
        if isinstance(item, tuple):
            text, opts = item
        else:
            text, opts = item, {}
        run = p.add_run()
        run.text = (bullet_char + "  " if not opts.get("nobullet") else "") + text
        run.font.name = font_name
        run.font.size = Pt(opts.get("size", font_size))
        run.font.bold = opts.get("bold", False)
        run.font.color.rgb = opts.get("color", color)
    return tb


def add_rect(slide, left, top, width, height, fill_color=LIGHT_BG, line_color=None):
    sh = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,
                                Inches(left), Inches(top), Inches(width), Inches(height))
    sh.fill.solid()
    sh.fill.fore_color.rgb = fill_color
    if line_color is None:
        sh.line.fill.background()
    else:
        sh.line.color.rgb = line_color
    sh.shadow.inherit = False
    return sh


def add_title_bar(slide, title_text, subtitle_text=None):
    # Top accent bar
    bar = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0), Inches(0), Inches(13.33), Inches(0.85))
    bar.fill.solid(); bar.fill.fore_color.rgb = NAVY
    bar.line.fill.background()
    add_textbox(slide, 0.5, 0.15, 12.3, 0.5, title_text,
                font_size=24, bold=True, color=RGBColor(0xFF, 0xFF, 0xFF))
    if subtitle_text:
        add_textbox(slide, 0.5, 0.50, 12.3, 0.4, subtitle_text,
                    font_size=12, color=RGBColor(0xC8, 0xD3, 0xE5))


def add_footer(slide, page_text):
    add_textbox(slide, 12.3, 7.05, 1.0, 0.3, page_text, font_size=10,
                color=RGBColor(0x90, 0x90, 0x90), align=PP_ALIGN.RIGHT)


# =========================================================
# Pre-build figures specific to deck
# =========================================================
def build_deck_figures():
    plt.rcParams.update({
        "font.family": "DejaVu Sans",
        "font.size": 11,
        "axes.spines.top": False,
        "axes.spines.right": False,
    })

    # --- Figure A: known vs holdout gap, 4 models × 2 datasets bar chart ---
    fig, ax = plt.subplots(figsize=(8, 4.0))
    models = ["Llama-3-8B", "Llama-3.1-8B", "Mistral-7B", "Qwen3-8B", "Gemma-2-9B"]
    cic_known   = [0.981, 0.976, 0.966, 0.974, 0.973]
    cic_holdout = [0.541, 0.568, 0.536, 0.706, 0.609]
    uns_known   = [0.984, 0.985, 0.979, 0.982, 0.975]
    uns_holdout = [0.989, 0.991, 0.971, 0.973, 0.989]
    x = np.arange(len(models))
    w = 0.18
    ax.bar(x - 1.5*w, cic_known,   w, color="#4C7BB6", label="CIC known")
    ax.bar(x - 0.5*w, cic_holdout, w, color="#E64D4D", label="CIC holdout")
    ax.bar(x + 0.5*w, uns_known,   w, color="#7AB8E8", label="UNSW known")
    ax.bar(x + 1.5*w, uns_holdout, w, color="#2EA85C", label="UNSW holdout")
    ax.set_xticks(x); ax.set_xticklabels(models, fontsize=9)
    ax.set_ylim(0.4, 1.05); ax.set_ylabel("Probe accuracy")
    ax.axhline(0.5, color="gray", lw=0.5, ls=":")
    ax.legend(loc="lower right", fontsize=8, ncol=2, frameon=False)
    ax.set_title("Probe accuracy: same probe, two stories", fontsize=12)
    plt.tight_layout()
    plt.savefig(DECK_FIGS / "deck_known_vs_holdout.png", dpi=180, bbox_inches="tight")
    plt.close(fig)

    # --- Figure B: V1 baseline comparison (gap bar) ---
    fig, ax = plt.subplots(figsize=(8, 3.5))
    methods = ["TF-IDF\nword 1-2g", "TF-IDF\nchar 3-5g", "Raw 79\nfeatures", "Llama-3.1\nLLM probe", "Qwen3\nLLM probe"]
    cic_gaps = [0.469, 0.319, 0.454, 0.408, 0.267]
    uns_gaps = [0.004, -0.003, 0.000, -0.006, 0.009]
    x = np.arange(len(methods))
    w = 0.36
    ax.bar(x - w/2, cic_gaps, w, color="#E64D4D", label="CIC2017 gap")
    ax.bar(x + w/2, uns_gaps, w, color="#2EA85C", label="UNSW gap")
    ax.set_xticks(x); ax.set_xticklabels(methods, fontsize=9)
    ax.set_ylabel("known acc − holdout acc")
    ax.axhline(0, color="black", lw=0.5)
    ax.legend(fontsize=9, frameon=False)
    ax.set_title("V1: simple baselines match LLM probe gap", fontsize=12)
    plt.tight_layout()
    plt.savefig(DECK_FIGS / "deck_v1_baseline.png", dpi=180, bbox_inches="tight")
    plt.close(fig)

    # --- Figure C: V2 top-axis cosines ---
    fig, ax = plt.subplots(figsize=(8, 4.0))
    pairs = ["Llama-3.1\nCICIDS", "Llama-3.1\nUNSW", "Qwen3\nCICIDS", "Qwen3\nUNSW",
             "Mistral\nCICIDS", "Mistral\nUNSW", "Gemma\nCICIDS", "Gemma\nUNSW"]
    top1_cos = [0.72, 0.70, 0.69, 0.50, 0.66, 0.50, 0.61, 0.13]
    feat = ["Init_Win_fwd", "dst_port", "Init_Win_fwd", "ct_dst_sport",
            "Init_Win_fwd", "ct_dst_sport", "Flow Pkts/s", "sttl"]
    colors = ["#E64D4D" if c >= 0.5 else "#7AB8E8" for c in top1_cos]
    ax.barh(np.arange(len(pairs)), top1_cos, color=colors, alpha=0.85)
    ax.set_yticks(np.arange(len(pairs))); ax.set_yticklabels(pairs, fontsize=9)
    ax.invert_yaxis()
    for i, (c, f) in enumerate(zip(top1_cos, feat)):
        ax.text(c + 0.01, i, f, va="center", fontsize=8.5)
    ax.axvline(0.5, color="gray", ls=":", lw=0.6)
    ax.text(0.51, len(pairs)-0.5, "single-axis\nthreshold (0.5)", fontsize=8, color="gray")
    ax.set_xlabel("|cos(probe direction, top-1 single-feature axis)|")
    ax.set_xlim(0, 1.0)
    ax.set_title("V2: in 7/8 cases probe direction ≈ one numeric field", fontsize=12)
    plt.tight_layout()
    plt.savefig(DECK_FIGS / "deck_v2_decomp.png", dpi=180, bbox_inches="tight")
    plt.close(fig)

    # --- Figure D: V3a direct-LLM heatmap (6 models × 5 datasets) ---
    fig, ax = plt.subplots(figsize=(8, 3.5))
    models = ["Llama-3-8B", "Llama-3.1-8B", "Mistral-7B", "Qwen3-8B-Inst", "Qwen3-8B-Base", "Gemma-2-9B"]
    datasets = ["CIC2017", "UNSW", "CIC2018", "IoT-23", "CTU-13"]
    M = np.array([
        [0.578, 0.569, 0.434, 0.545, 0.554],
        [0.605, 0.557, 0.437, 0.544, 0.540],
        [0.500, 0.369, 0.537, 0.617, 0.582],
        [0.516, 0.415, 0.542, 0.349, 0.481],
        [0.565, 0.330, 0.461, 0.375, 0.511],
        [0.555, 0.438, 0.455, 0.483, 0.595],
    ])
    im = ax.imshow(M, cmap="RdYlGn", vmin=0.3, vmax=0.7, aspect="auto")
    ax.set_xticks(range(len(datasets))); ax.set_xticklabels(datasets, fontsize=9)
    ax.set_yticks(range(len(models))); ax.set_yticklabels(models, fontsize=9)
    for i in range(len(models)):
        for j in range(len(datasets)):
            ax.text(j, i, f"{M[i,j]:.2f}", ha="center", va="center",
                    color="black", fontsize=9)
    plt.colorbar(im, ax=ax, label="AUROC", fraction=0.04, pad=0.02)
    ax.set_title("V3a: LLM zero-shot direct yes/no — all near chance (0.5)",
                 fontsize=11)
    plt.tight_layout()
    plt.savefig(DECK_FIGS / "deck_v3a_direct.png", dpi=180, bbox_inches="tight")
    plt.close(fig)

    # --- Figure E: V3b steering bar (cic_attack vs harmful vs random) ---
    fig, ax = plt.subplots(figsize=(8, 3.5))
    layers = [9, 18, 26, 31]
    atk_spans  = [5.69, 5.60, 3.86, 4.00]
    harm_spans = [6.42, 9.75, 4.75, 9.44]
    rand_spans = [5.07, 5.38, 6.53, 8.19]
    x = np.arange(len(layers)); w = 0.27
    ax.bar(x - w, atk_spans,  w, color="#E64D4D", label="cic_attack")
    ax.bar(x,     harm_spans, w, color="#2EA85C", label="harmful (JBB)")
    ax.bar(x + w, rand_spans, w, color="#999999", label="random control")
    ax.set_xticks(x); ax.set_xticklabels([f"L={l}" for l in layers])
    ax.set_ylabel("Steering span (Δ log p(yes/no))")
    ax.legend(fontsize=9, frameon=False)
    ax.set_title("V3b steering span (Qwen3-Base): attack ≈ random; harmful > random",
                 fontsize=11)
    plt.tight_layout()
    plt.savefig(DECK_FIGS / "deck_v3b_steering.png", dpi=180, bbox_inches="tight")
    plt.close(fig)

    # --- Figure F: V4 cross-dataset transfer heatmap ---
    fig, ax = plt.subplots(figsize=(8, 3.0))
    models = ["Llama-3-8B", "Mistral-7B", "Qwen3-8B-Base", "Gemma-2-9B"]
    targets = ["CIC2018\n(same schema)", "IoT-23\n(Zeek)", "CTU-13\n(NetFlow)"]
    M = np.array([
        [0.503, 0.605, 0.599],
        [0.371, 0.630, 0.651],
        [0.424, 0.646, 0.679],
        [0.429, 0.658, 0.679],
    ])
    im = ax.imshow(M, cmap="RdYlGn", vmin=0.35, vmax=0.85, aspect="auto")
    ax.set_xticks(range(len(targets))); ax.set_xticklabels(targets, fontsize=9)
    ax.set_yticks(range(len(models))); ax.set_yticklabels(models, fontsize=9)
    for i in range(len(models)):
        for j in range(len(targets)):
            ax.text(j, i, f"{M[i,j]:.2f}", ha="center", va="center", fontsize=10)
    plt.colorbar(im, ax=ax, label="AUROC", fraction=0.04, pad=0.02)
    ax.set_title("V4: zero-shot transfer of CIC2017 attack-direction to 3 unseen datasets",
                 fontsize=11)
    plt.tight_layout()
    plt.savefig(DECK_FIGS / "deck_v4_transfer.png", dpi=180, bbox_inches="tight")
    plt.close(fig)

    # --- Figure G: per-attack-type comparison CIC train vs holdout ---
    fig, axes = plt.subplots(1, 2, figsize=(9, 3.2), sharey=False)
    # CICIDS
    types_cic = ["DoS Hulk", "PortScan", "DDoS", "Brute Force", "Bot\n(holdout)", "Heartbleed\n(holdout)", "Infilt.\n(holdout)"]
    rates_cic = [6.4e6, 422, 1.0e6, 645, 9.6, 6.3e4, 10.1]
    is_holdout = [False, False, False, False, True, True, True]
    cmap = ["#7AB8E8" if not h else "#E64D4D" for h in is_holdout]
    axes[0].bar(types_cic, np.log10(np.array(rates_cic) + 1), color=cmap)
    axes[0].set_yticks([0, 1, 2, 3, 4, 5, 6, 7])
    axes[0].set_yticklabels(["1", "10", "100", "1k", "10k", "100k", "1M", "10M"])
    axes[0].set_ylabel("Flow Bytes/s (log scale)")
    axes[0].set_title("CIC2017: train vs holdout differ ×10⁵", fontsize=11)
    axes[0].tick_params(axis="x", labelsize=8)
    plt.setp(axes[0].get_xticklabels(), rotation=20, ha="right")

    # UNSW (use ttl as a categorical-like)
    types_uns = ["Exploits", "Recon.", "Generic", "Fuzzers", "Dos", "Backdoor\n(holdout)", "Shellcode\n(holdout)", "Worms\n(holdout)"]
    ttl_uns = [254, 254, 254, 254, 254, 254, 254, 254]
    holdout_uns = [False, False, False, False, False, True, True, True]
    cmap_u = ["#7AB8E8" if not h else "#E64D4D" for h in holdout_uns]
    axes[1].bar(types_uns, ttl_uns, color=cmap_u)
    axes[1].set_ylabel("source TTL")
    axes[1].set_ylim(0, 280)
    axes[1].set_title("UNSW: all attacks share TTL=254, service=None, FIN", fontsize=11)
    axes[1].tick_params(axis="x", labelsize=8)
    plt.setp(axes[1].get_xticklabels(), rotation=20, ha="right")
    axes[1].axhline(31, color="gray", lw=0.6, ls=":")
    axes[1].text(7.5, 35, "normal TTL ≈ 31", fontsize=8, color="gray", ha="right")

    plt.tight_layout()
    plt.savefig(DECK_FIGS / "deck_pattern_explanation.png", dpi=180, bbox_inches="tight")
    plt.close(fig)


# =========================================================
# Build slides
# =========================================================
def build_deck():
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)
    blank = prs.slide_layouts[6]   # blank layout

    # ============ Slide 1: title ============
    s = prs.slides.add_slide(blank)
    # Background gradient strip
    bar = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0), Inches(0), Inches(13.33), Inches(7.5))
    bar.fill.solid(); bar.fill.fore_color.rgb = NAVY
    bar.line.fill.background()
    add_textbox(s, 1.0, 2.4, 11.3, 1.4,
                "When 99 % Probe Accuracy Lies",
                font_size=44, bold=True, color=RGBColor(0xFF, 0xFF, 0xFF))
    add_textbox(s, 1.0, 3.6, 11.3, 0.9,
                "An LLM does not have a transferable “network-attack” concept",
                font_size=22, color=RGBColor(0xC8, 0xD3, 0xE5))
    # Author + course
    add_textbox(s, 1.0, 5.5, 6.0, 0.4,
                "Yihuai Hong  ·  New York University",
                font_size=16, color=RGBColor(0xE8, 0xEE, 0xF7))
    add_textbox(s, 1.0, 6.0, 6.0, 0.4,
                "End-of-term project presentation, 2025",
                font_size=14, color=RGBColor(0xB8, 0xC4, 0xD8))

    # ============ Slide 2 (NEW): Background — refusal direction in LLM safety ============
    s = prs.slides.add_slide(blank)
    add_title_bar(s, "Background — concepts as linear directions in LLM hidden space",
                  "Prior LLM-safety work shows harmful vs benign instructions are linearly separable")

    # Left: prior findings
    add_textbox(s, 0.5, 1.2, 6.0, 0.5, "What prior work has shown",
                font_size=18, bold=True, color=NAVY)
    add_bullets(s, 0.5, 1.75, 6.0, 5.0,
                [("Linear separability:", {"size": 14, "bold": True, "color": NAVY, "nobullet": True}),
                 "A single linear probe on residual-stream activations distinguishes harmful prompts (e.g. ‘how to build a bomb’) from benign prompts (e.g. ‘how to bake a cake’) at  >  98 % accuracy.",
                 ("", {"nobullet": True}),
                 ("Causal effect — the ‘refusal direction’:", {"size": 14, "bold": True, "color": NAVY, "nobullet": True}),
                 "Adding ±α · d_harmful to the residual stream flips the model between refusal and compliance.",
                 "Zeroing-out d_harmful jailbreaks the safety-tuned model.",
                 ("", {"nobullet": True}),
                 ("Generality:", {"size": 14, "bold": True, "color": NAVY, "nobullet": True}),
                 "Single direction works across many prompt formats and harm categories.",
                 "Robust across model families (Llama-2/3, Qwen, Mistral, Gemma).",
                 ], font_size=12)

    # Right: key citations + the natural extension
    add_rect(s, 6.7, 1.2, 6.2, 2.6, fill_color=LIGHT_BG)
    add_textbox(s, 6.9, 1.32, 5.8, 0.4, "Key references",
                font_size=14, bold=True, color=NAVY)
    add_bullets(s, 6.9, 1.75, 5.8, 1.9,
                ["Arditi et al. (NeurIPS 2024) — ‘Refusal in LLMs is mediated by a single direction’",
                 "Zou et al. (2023) — ‘Representation Engineering’",
                 "Li et al. (NeurIPS 2024) — ‘Inference-Time Intervention for truthfulness’",
                 "Chao et al. (NeurIPS 2024) — JailbreakBench dataset"],
                font_size=11, color=GRAY_T)

    # Bottom-right: the leap
    add_rect(s, 6.7, 4.0, 6.2, 2.9, fill_color=RGBColor(0xFF, 0xF5, 0xF0), line_color=ACCENT)
    add_textbox(s, 6.9, 4.15, 5.8, 0.5, "The natural question for this work",
                font_size=15, bold=True, color=ACCENT)
    add_textbox(s, 6.9, 4.65, 5.8, 2.2,
                "If a single direction encodes ‘harmful natural-language instructions,’\n"
                "does a similar direction encode ‘malicious network traffic’?\n\n"
                "We attempted to extract such a direction on CICIDS2017 and UNSW-NB15.\n"
                "The naive probing pipeline reports >99 % CV AUROC.\n\n"
                "This talk:  why that 99 % is misleading.",
                font_size=12, color=GRAY_T)

    add_footer(s, "2 / 16")

    # ============ Slide 3: the puzzle ============
    s = prs.slides.add_slide(blank)
    add_title_bar(s, "The puzzle: a probe at 99 % accuracy",
                  "But same probe, same model, two completely different stories")
    add_bullets(s, 0.6, 1.2, 7.0, 4.0,
                [("Standard probing recipe says:", {"size": 18, "bold": True, "color": NAVY}),
                 "Take an LLM, format network-flow logs as text",
                 "Train a logistic probe on hidden states  →  >99 % CV AUROC",
                 "Conclusion (per the recipe):  the LLM has an internal  “attack”  concept",
                 ("", {"nobullet": True}),
                 ("BUT, when we hold out attack types unseen at training:", {"size": 18, "bold": True, "color": ACCENT}),
                 "UNSW-NB15:  98 % known →  99 % zero-day  (excellent)",
                 "CICIDS-2017: 98 % known →  54 % zero-day  (collapse)",
                 ("Either the LLM has UNSW concept but no CIC concept …", {"nobullet": True, "size": 14, "color": GRAY_T}),
                 ("…or the high probe accuracy was misleading from the start.", {"nobullet": True, "size": 14, "color": ACCENT, "bold": True}),
                ], font_size=15)
    s.shapes.add_picture(str(DECK_FIGS / "deck_known_vs_holdout.png"),
                          Inches(7.7), Inches(1.5), width=Inches(5.4))
    add_footer(s, "3 / 16")

    # ============ Slide 3: setup ============
    s = prs.slides.add_slide(blank)
    add_title_bar(s, "Setup",
                  "6 models, 5 datasets, and a probe-validation protocol")
    # Models panel
    add_rect(s, 0.5, 1.2, 6.0, 2.6)
    add_textbox(s, 0.7, 1.3, 5.6, 0.4, "Models", font_size=16, bold=True, color=NAVY)
    add_bullets(s, 0.7, 1.75, 5.6, 2.0,
                ["Meta-Llama-3-8B-Instruct",
                 "Llama-3.1-8B-Instruct",
                 "Mistral-7B-Instruct-v0.3",
                 "Qwen3-8B-Instruct  &  Qwen3-8B-Base",
                 "Gemma-2-9B-it"], font_size=14)
    # Datasets panel
    add_rect(s, 6.7, 1.2, 6.1, 2.6)
    add_textbox(s, 6.9, 1.3, 5.7, 0.4, "Datasets", font_size=16, bold=True, color=NAVY)
    add_bullets(s, 6.9, 1.75, 5.7, 2.0,
                [("Train + holdout test:", {"nobullet": True, "size": 13, "bold": True}),
                 "CICIDS2017 (CICFlowMeter, 18 num. fields)",
                 "UNSW-NB15 (Argus, 19 fields incl. service / state)",
                 ("Zero-shot transfer test (3 unseen datasets):", {"nobullet": True, "size": 13, "bold": True}),
                 "CICIDS2018 (same schema)  ·  IoT-23 (Zeek)  ·  CTU-13 (NetFlow)"],
                font_size=13)
    # Pipeline diagram
    pipeline_y = 4.3
    add_textbox(s, 0.5, pipeline_y, 12.3, 0.4, "Pipeline (extract reps then probe):",
                font_size=14, bold=True, color=NAVY)
    box_w = 2.4
    boxes = ["Network log  →  text", "LLM forward pass", "last-token  hidden state",
             "logistic probe", "known/holdout AUROC"]
    for i, txt in enumerate(boxes):
        x0 = 0.5 + i * (box_w + 0.05)
        if i % 2 == 0:
            color = LIGHT_BG
        else:
            color = RGBColor(0xE3, 0xEB, 0xF7)
        add_rect(s, x0, pipeline_y + 0.6, box_w, 0.7, fill_color=color)
        add_textbox(s, x0, pipeline_y + 0.65, box_w, 0.6, txt, font_size=11,
                    bold=True, align=PP_ALIGN.CENTER, color=NAVY)
        if i < len(boxes) - 1:
            add_textbox(s, x0 + box_w - 0.18, pipeline_y + 0.78, 0.45, 0.3, "→",
                        font_size=20, color=NAVY, bold=True)
    add_textbox(s, 0.5, pipeline_y + 1.5, 12.3, 0.6,
                "Holdout split: 3 attack types removed from train (CIC: Bot/Heartbleed/Infiltration; UNSW: Backdoor/Shellcode/Worms)",
                font_size=13, color=GRAY_T)
    add_footer(s, "4 / 16")

    # ============ Slide 3.5 (NEW): holdout split design ============
    s = prs.slides.add_slide(blank)
    add_title_bar(s, "Holdout split design",
                  "We hold out 3 attack types entirely from training; the probe sees them only at test time")

    # CICIDS panel
    add_rect(s, 0.4, 1.2, 6.3, 5.4)
    add_textbox(s, 0.6, 1.32, 6.0, 0.5, "CICIDS2017", font_size=20, bold=True, color=NAVY)
    add_textbox(s, 0.6, 1.85, 6.0, 0.4, "Train set (probe sees these)",
                font_size=14, bold=True, color=RGBColor(0x4C, 0x7B, 0xB6))
    add_bullets(s, 0.6, 2.25, 6.0, 2.4,
                ["DoS Hulk     (n = 230 k)",
                 "PortScan     (n = 159 k)",
                 "DDoS         (n = 128 k)",
                 "DoS GoldenEye, slowloris, Slowhttptest",
                 "FTP-Patator  ·  SSH-Patator",
                 "Web Attack — BruteForce  ·  XSS"],
                font_size=13, font_name="Consolas",
                color=RGBColor(0x4C, 0x7B, 0xB6))
    add_textbox(s, 0.6, 4.55, 6.0, 0.4, "Holdout (zero-day eval)",
                font_size=14, bold=True, color=ACCENT)
    add_bullets(s, 0.6, 4.95, 6.0, 1.5,
                ["Bot          (n = 1 956)",
                 "Heartbleed   (n = 11)",
                 "Infiltration (n = 36)"],
                font_size=13, font_name="Consolas", color=ACCENT)

    # UNSW panel
    add_rect(s, 6.95, 1.2, 6.0, 5.4)
    add_textbox(s, 7.15, 1.32, 5.7, 0.5, "UNSW-NB15", font_size=20, bold=True, color=NAVY)
    add_textbox(s, 7.15, 1.85, 5.7, 0.4, "Train set",
                font_size=14, bold=True, color=RGBColor(0x4C, 0x7B, 0xB6))
    add_bullets(s, 7.15, 2.25, 5.7, 2.4,
                ["Exploits     (n = 27 k)",
                 "Generic      (n = 25 k)",
                 "Fuzzers      (n = 22 k)",
                 "Reconnaissance (n = 13 k)",
                 "DoS          (n = 17 k)",
                 "Analysis"],
                font_size=13, font_name="Consolas",
                color=RGBColor(0x4C, 0x7B, 0xB6))
    add_textbox(s, 7.15, 4.55, 5.7, 0.4, "Holdout (zero-day eval)",
                font_size=14, bold=True, color=ACCENT)
    add_bullets(s, 7.15, 4.95, 5.7, 1.5,
                ["Backdoor    (n = 1 983)",
                 "Shellcode   (n = 1 511)",
                 "Worms       (n = 171)"],
                font_size=13, font_name="Consolas", color=ACCENT)

    add_textbox(s, 0.4, 6.78, 12.5, 0.5,
                "Probe is trained ONLY on blue attacks + Normal.  Red attacks are pure zero-day.  Same Normal pool is split 60/20/20 across train/known-test/holdout-test.",
                font_size=12, color=GRAY_T, align=PP_ALIGN.CENTER)
    add_footer(s, "5 / 16")

    # ============ Slide 4: sample data shown to LLM ============
    s = prs.slides.add_slide(blank)
    add_title_bar(s, "What does the LLM actually see?",
                  "One representative formatted log per dataset (truncated to fit)")

    samples = [
        ("CICIDS2017 — Bot (holdout)",
         "protocol: tcp\nsource port: 42544\ndestination port: 8080\nflow duration: 60 202 640\nflow bytes per second: 9.6\nflow packets per second: 0.299\nSYN/RST/ACK: 0/0/0   PSH: 1\ninitial forward window bytes: 29 200",
         ACCENT),
        ("UNSW-NB15 — Backdoor (holdout)",
         "protocol: tcp\nsource port: 29 335\ndestination port: 514\nduration: 1.29 s\nsource TTL: 254  /  destination TTL: 252\nservice: None     connection state: FIN\nsource packets: 10  /  destination packets: 8\nct_srv_src: 4   ct_dst_sport_ltm: 1",
         ACCENT),
        ("CICIDS2018 — DDOS-HOIC (transfer)",
         "protocol: 6\ndestination port: 80\nflow duration: 2 316\ntotal forward packets: 2   backward: 0\nflow bytes per second: 0.0\nflow packets per second: 863.6\nSYN/RST/ACK: 0/0/1   PSH: 0\ninitial forward window bytes: 32 738",
         NAVY),
        ("IoT-23 — PortScan (transfer)",
         "protocol: tcp\nsource port: 43 746\ndestination port: 80\nconnection state: S0\nhistory: S\noriginator packets: 1\nresponder packets: 0\noriginator IP bytes: 40\nresponder IP bytes: 0",
         NAVY),
        ("CTU-13 — Botnet (transfer)",
         "source port: 2 593\ndestination port: 443\nduration: 1.10 s\ntotal packets: 16   total bytes: 5 565\nsource bytes: 1 420   source ToS: 0\nbytes per packet: 347.8\npacket/byte ratio: 0.0029",
         NAVY),
    ]
    panel_w = 4.1; panel_h = 2.7; row_y = [1.2, 4.05]
    for i, (title, body, color) in enumerate(samples):
        r = i // 3; c = i % 3
        x0 = 0.4 + c * (panel_w + 0.15)
        y0 = row_y[r]
        add_rect(s, x0, y0, panel_w, panel_h, fill_color=LIGHT_BG, line_color=color)
        add_textbox(s, x0 + 0.15, y0 + 0.08, panel_w - 0.3, 0.4, title,
                    font_size=12, bold=True, color=color)
        add_textbox(s, x0 + 0.15, y0 + 0.55, panel_w - 0.3, panel_h - 0.6,
                    body, font_size=9.5, color=GRAY_T, font_name="Consolas")
    add_textbox(s, 0.4, 6.85, 12.5, 0.5,
                "Color = role in our protocol. Red boxes are attacks the train probe never saw. "
                "Blue boxes are completely separate datasets used for V4 transfer.",
                font_size=12, color=GRAY_T, align=PP_ALIGN.CENTER)
    add_footer(s, "6 / 16")

    # ============ Slide 5: validation protocol V1-V4 ============
    s = prs.slides.add_slide(blank)
    add_title_bar(s, "Our methodological contribution",
                  "Probe accuracy alone is unreliable.  We propose 4 validation tests.")
    cells = [
        ("V1", "Baseline equivalence",
         "Does TF-IDF or raw-feature LR achieve the same probe gap as the LLM?",
         "If YES (LLM ≈ baseline)  →  LLM has no value-add."),
        ("V2", "Decomposition",
         "Is cos(direction, some single-feature axis) > 0.5  for any input field?",
         "If YES (some single field aligns)  →  direction is a single-feature shortcut, not a concept."),
        ("V3", "Causal effect",
         "Does the model itself use this direction?  (zero-shot yes/no + activation steering)",
         "If NO (model can't use it)  →  direction is a passive feature, not a decision axis."),
        ("V4", "Transfer",
         "Does the direction separate attack vs normal in a previously-unseen dataset?",
         "If NO (transfer at chance)  →  direction is a dataset-specific hash."),
    ]
    for i, (key, name, desc, verdict) in enumerate(cells):
        x0 = 0.5 + (i % 2) * 6.4
        y0 = 1.2 + (i // 2) * 2.8
        add_rect(s, x0, y0, 6.1, 2.5)
        add_textbox(s, x0 + 0.2, y0 + 0.15, 1.0, 0.55, key,
                    font_size=32, bold=True, color=NAVY)
        add_textbox(s, x0 + 1.3, y0 + 0.2, 4.6, 0.45, name,
                    font_size=18, bold=True, color=NAVY)
        add_textbox(s, x0 + 0.2, y0 + 0.85, 5.7, 1.0, desc,
                    font_size=14, color=GRAY_T)
        add_textbox(s, x0 + 0.2, y0 + 1.75, 5.7, 0.7,
                    verdict,
                    font_size=11, bold=True, color=ACCENT)
    add_textbox(s, 0.5, 6.95, 12.3, 0.4,
                "A real concept passes V1 (no), V2 (no), V3 (yes), V4 (yes).  Our attack direction fails all four.",
                font_size=14, bold=True, color=ACCENT, align=PP_ALIGN.CENTER)
    add_footer(s, "7 / 16")

    # ============ Slide 5: V1 baseline equivalence ============
    s = prs.slides.add_slide(blank)
    add_title_bar(s, "V1 — Bag-of-features baselines reproduce the gap",
                  "TF-IDF on text and even raw numeric features match the LLM probe")
    s.shapes.add_picture(str(DECK_FIGS / "deck_v1_baseline.png"),
                          Inches(0.4), Inches(1.2), width=Inches(7.6))
    add_bullets(s, 8.2, 1.4, 4.9, 5.5,
                [("Headline:", {"size": 16, "bold": True, "color": NAVY, "nobullet": True}),
                 ("LR on 79 raw numeric fields", {"size": 14, "bold": True, "color": ACCENT}),
                 "→ CIC gap +0.45,  UNSW gap 0.00",
                 ("≡ Llama-3 LLM probe", {"size": 14, "bold": True, "color": ACCENT}),
                 ("", {"nobullet": True}),
                 ("Implication:", {"size": 16, "bold": True, "color": NAVY, "nobullet": True}),
                 "8B parameters of LLM contribute nothing beyond the surface lexical / numeric statistics on this task.",
                 ("", {"nobullet": True}),
                 ("Caveat:", {"size": 14, "bold": True, "color": GRAY_T, "nobullet": True}),
                 "Qwen3 has a slightly smaller CIC gap (0.27); we examine why later."],
                font_size=13)
    add_footer(s, "8 / 16")

    # ============ Slide 6: V2 decomposition ============
    s = prs.slides.add_slide(blank)
    add_title_bar(s, "V2 — The probe direction is one or two numeric fields",
                  "Build a unit axis from train normals sorted by each field, take cos with the probe direction")
    s.shapes.add_picture(str(DECK_FIGS / "deck_v2_decomp.png"),
                          Inches(0.4), Inches(1.2), width=Inches(7.5))
    add_bullets(s, 8.0, 1.4, 5.1, 5.0,
                [("In 7 / 8 (model, dataset) cases:", {"size": 16, "bold": True, "color": NAVY, "nobullet": True}),
                 "top-1 field axis explains ≥ 0.5 of the direction",
                 "Init_Win_bytes_fwd dominates 4 cases",
                 "destination_port dominates 2 cases",
                 ("", {"nobullet": True}),
                 ("Adversarial check (P9):", {"size": 14, "bold": True, "color": NAVY, "nobullet": True}),
                 "perturb 5 candidate fields by ×{0, 0.01, 100} —  only Init_Win_bytes_fwd flips the probe.",
                 ("", {"nobullet": True}),
                 ("Lone exception:", {"size": 14, "bold": True, "color": GREEN, "nobullet": True}),
                 "Qwen3 / UNSW (cos = 0.50, distributed). Still fails V4."],
                font_size=12)
    add_footer(s, "9 / 16")

    # ============ Slide 7: V3a direct LLM ============
    s = prs.slides.add_slide(blank)
    add_title_bar(s, "V3a — Asked directly, the LLM is at chance",
                  '"You are a network security analyst.  Is this connection malicious?  yes / no"')
    s.shapes.add_picture(str(DECK_FIGS / "deck_v3a_direct.png"),
                          Inches(0.4), Inches(1.2), width=Inches(8.4))
    add_bullets(s, 9.0, 1.4, 4.0, 5.0,
                [("All 30 cells in 0.33 — 0.62", {"size": 14, "bold": True, "color": ACCENT, "nobullet": True}),
                 ("(random ± 0.13)", {"nobullet": True, "size": 12, "color": GRAY_T}),
                 ("", {"nobullet": True}),
                 ("Probe AUROC on the same hidden states:", {"size": 13, "bold": True, "color": NAVY, "nobullet": True}),
                 ("≥ 0.99 in-domain", {"size": 13, "color": NAVY, "nobullet": True}),
                 ("", {"nobullet": True}),
                 ("Gap of 35–60 AUROC points", {"size": 14, "bold": True, "color": ACCENT, "nobullet": True}),
                 ("between ‘probe extracts’ and ‘model uses’.", {"nobullet": True, "size": 12}),
                 ("", {"nobullet": True}),
                 ("Information is in the hidden state.", {"nobullet": True, "size": 13, "bold": True, "color": NAVY}),
                 ("The model does not use it.", {"nobullet": True, "size": 13, "bold": True, "color": ACCENT}),
                 ], font_size=13)
    add_footer(s, "10 / 16")

    # ============ Slide 8: V3b steering ============
    s = prs.slides.add_slide(blank)
    add_title_bar(s, "V3b — Steering the residual stream",
                  "Add α · direction at multiple layers (CAA-style) and measure Δ log p(yes) / log p(no)")
    s.shapes.add_picture(str(DECK_FIGS / "deck_v3b_steering.png"),
                          Inches(0.4), Inches(1.2), width=Inches(8.0))
    add_bullets(s, 8.7, 1.4, 4.4, 5.5,
                [("Attack direction span:", {"size": 14, "bold": True, "color": ACCENT, "nobullet": True}),
                 ("≈ random-direction span at deep layers", {"nobullet": True, "size": 13}),
                 ("", {"nobullet": True}),
                 ("Harmful direction (JBB):", {"size": 14, "bold": True, "color": GREEN, "nobullet": True}),
                 ("1.4-1.8× random, monotonic in α", {"nobullet": True, "size": 13}),
                 ("(positive control passes V3)", {"nobullet": True, "size": 12, "color": GRAY_T}),
                 ("", {"nobullet": True}),
                 ("Layer nuance:", {"size": 13, "bold": True, "color": NAVY, "nobullet": True}),
                 "At very early layers (L=4) Qwen3 attack direction has 4× random span — but still fails V1, V2, V4."
                 ], font_size=13)
    add_footer(s, "11 / 16")

    # ============ Slide 9: V4 transfer ============
    s = prs.slides.add_slide(blank)
    add_title_bar(s, "V4 — Cross-dataset transfer fails decisively",
                  "Project new-dataset hidden states through the CICIDS2017 attack-direction.  AUROC vs the new dataset's labels")
    s.shapes.add_picture(str(DECK_FIGS / "deck_v4_transfer.png"),
                          Inches(0.4), Inches(1.2), width=Inches(7.8))
    add_bullets(s, 8.4, 1.4, 4.7, 5.5,
                [("In-domain probe ≥ 0.90 on every model × dataset pair.", {"size": 13, "color": NAVY, "nobullet": True}),
                 ("", {"nobullet": True}),
                 ("Maximum cross-dataset AUROC = 0.679", {"size": 14, "bold": True, "color": ACCENT, "nobullet": True}),
                 ("(Qwen3-Base → CTU-13)", {"nobullet": True, "size": 12}),
                 ("", {"nobullet": True}),
                 ("Same-schema (CIC2017 → CIC2018) is below chance for 3 / 4 models", {"size": 14, "bold": True, "color": ACCENT, "nobullet": True}),
                 ("    The CIC2017 direction is overfit to its specific attack mix and actively misclassifies CIC2018 attacks.", {"nobullet": True, "size": 11}),
                 ("", {"nobullet": True}),
                 ("Robustness: same conclusion holds for the full LR probe (P3 experiment).", {"nobullet": True, "size": 12, "color": GRAY_T}),
                 ], font_size=13)
    add_footer(s, "12 / 16")

    # ============ Slide 10: positive control ============
    s = prs.slides.add_slide(blank)
    add_title_bar(s, "Positive control — the harmful direction passes V3 + V4",
                  "Apply the same protocol to a known-real concept to validate the protocol itself")
    add_bullets(s, 0.5, 1.2, 6.2, 5.5,
                [("Two different ‘attacks’ in this work:", {"size": 16, "bold": True, "color": NAVY, "nobullet": True}),
                 ("(A) Network attack  =  malicious flow record  ←  main task", {"nobullet": True, "size": 13, "bold": True, "color": ACCENT}),
                 ("(B) Harmful instruction  =  jailbreak prompt  ←  positive control", {"nobullet": True, "size": 13, "bold": True, "color": GREEN}),
                 ("", {"nobullet": True}),
                 ("Concept (B) source — JBB (JailbreakBench):", {"size": 14, "bold": True, "color": NAVY, "nobullet": True}),
                 "100 harmful prompts (e.g. ‘write a tutorial on how to make a bomb’)",
                 "100 benign prompts (matched topics, no harmful intent)",
                 "Direction:  mean(harmful reps) − mean(benign reps)",
                 ("", {"nobullet": True}),
                 ("Reference literature:  Arditi et al. (NeurIPS 2024) ‘Refusal Direction’.", {"nobullet": True, "size": 12, "color": GRAY_T}),
                 ], font_size=13)
    # Right side: cosine table
    add_rect(s, 6.8, 1.3, 6.1, 5.6)
    add_textbox(s, 7.0, 1.4, 5.7, 0.4,
                "cos(harmful direction,  attack direction)",
                font_size=14, bold=True, color=NAVY)
    add_textbox(s, 7.0, 1.85, 5.7, 0.35,
                "raw-text JBB extraction, max over layers",
                font_size=10, color=GRAY_T)
    rows = [("Llama-3-8B",  "0.02 / 0.14"),
            ("Mistral-7B",  "0.13 / 0.14"),
            ("Qwen3-8B",    "0.30 / 0.47"),
            ("Gemma-2-9B",  "0.13 / 0.18")]
    add_textbox(s, 7.0, 2.4, 3.0, 0.4, "model", font_size=12, bold=True, color=NAVY)
    add_textbox(s, 10.0, 2.4, 3.0, 0.4, "CICIDS  /  UNSW", font_size=12, bold=True, color=NAVY)
    for i, (m, c) in enumerate(rows):
        add_textbox(s, 7.0, 2.85 + i*0.5, 3.0, 0.4, m, font_size=13, color=GRAY_T)
        add_textbox(s, 10.0, 2.85 + i*0.5, 3.0, 0.4, c, font_size=13, color=GRAY_T)
    add_textbox(s, 7.0, 5.0, 5.7, 0.4,
                "Verdict on harmful direction:",
                font_size=14, bold=True, color=GREEN)
    add_bullets(s, 7.0, 5.4, 5.7, 1.7,
                ["V3 steering: 1.4-1.8 × random, monotonic in α  (✓)",
                 "V4 cross-prompt transfer: high AUROC across formats  (✓)",
                 "Even at peak cos 0.47 (Qwen3 / UNSW), netsec dir fails V4 — adjacent ≠ same concept."],
                font_size=11, color=GRAY_T)
    add_footer(s, "13 / 16")

    # ============ Slide 11: why CIC fails / UNSW doesn't ============
    s = prs.slides.add_slide(blank)
    add_title_bar(s, "Why CIC fails but UNSW doesn't",
                  "The 'good' generalization on UNSW is a holdout-design artifact, not LLM understanding")
    s.shapes.add_picture(str(DECK_FIGS / "deck_pattern_explanation.png"),
                          Inches(0.4), Inches(1.2), width=Inches(8.5))
    add_bullets(s, 9.1, 1.4, 4.0, 5.5,
                [("CIC train mix:", {"size": 14, "bold": True, "color": NAVY, "nobullet": True}),
                 "94 % DoS / DDoS / scan",
                 "high-rate signature",
                 ("", {"nobullet": True}),
                 ("CIC holdout:", {"size": 14, "bold": True, "color": ACCENT, "nobullet": True}),
                 "97 % Bot (low-rate ≈ Normal)",
                 "→ different surface signature",
                 ("", {"nobullet": True}),
                 ("UNSW (both train + holdout):", {"size": 14, "bold": True, "color": GREEN, "nobullet": True}),
                 "TTL = 254 + service = None + state = FIN",
                 "Same lexical fingerprint",
                 "→ TF-IDF char-ngrams already at 0.99 holdout"
                 ], font_size=12)
    add_footer(s, "14 / 16")

    # ============ Slide 12: things that reduce gap but don't unlock concept ============
    s = prs.slides.add_slide(blank)
    add_title_bar(s, "What can reduce the CIC gap (but does not unlock a concept)",
                  "Three independent interventions reduce gap.  None of them passes V1, V2, or V4.")
    rows = [
        ("key_value prompt format",          "0.40  →  0.24",  "Less format-specific overfit"),
        ("PCA-LR probe (top 128 components)","0.50  →  0.17",  "Throws away high-rank overfit modes"),
        ("Joint train (CIC + UNSW)",         "0.27  →  0.13",  "Larger, more diverse train mix"),
    ]
    add_textbox(s, 0.5, 1.3, 12.3, 0.4,
                "intervention                                         CIC gap  reduction                   why it helps",
                font_size=12, bold=True, color=NAVY)
    for i, (a, b, c) in enumerate(rows):
        y0 = 1.9 + i * 0.95
        add_rect(s, 0.5, y0, 12.3, 0.85)
        add_textbox(s, 0.7,  y0 + 0.2, 4.5, 0.5, a, font_size=15, bold=True, color=NAVY)
        add_textbox(s, 5.6,  y0 + 0.2, 3.5, 0.5, b, font_size=16, bold=True, color=ACCENT, align=PP_ALIGN.CENTER)
        add_textbox(s, 9.2,  y0 + 0.2, 3.5, 0.5, c, font_size=13, color=GRAY_T)

    add_textbox(s, 0.5, 5.0, 12.3, 1.6,
                "Each of the three brings the LLM probe closer to the TF-IDF baseline gap, but none unlocks  V4 transfer.\n"
                "These are  ways of making the probe  less brittle, not evidence of a hidden  concept.\n"
                "A concept would also need to (a) outperform TF-IDF, (b) be steerable, (c) transfer cross-dataset.",
                font_size=14, color=GRAY_T)
    add_footer(s, "15 / 16")

    # ============ Slide 13: conclusions ============
    s = prs.slides.add_slide(blank)
    add_title_bar(s, "Conclusions",
                  "When 99 % probe accuracy lies, what should we do?")
    add_bullets(s, 0.5, 1.2, 12.3, 4.5,
                [("1.  An LLM probe at 99 % accuracy is not evidence of an internal concept.", {"size": 16, "bold": True, "color": NAVY, "nobullet": True}),
                 "We propose a 4-test protocol (V1 baseline, V2 decomposition, V3 causal, V4 transfer).",
                 ("", {"nobullet": True}),
                 ("2.  On 6 LLMs × 5 network-flow datasets, the attack direction fails all 4.", {"size": 16, "bold": True, "color": ACCENT, "nobullet": True}),
                 "TF-IDF baseline matches the probe.  Direction reduces to 1-2 numeric fields.",
                 "Direct LLM yes / no is at chance.  Steering is causally inert.",
                 "Cross-dataset transfer never exceeds AUROC 0.68;  same-schema is below chance.",
                 ("", {"nobullet": True}),
                 ("3.  The same protocol  certifies the JBB harmful direction as a real concept.", {"size": 16, "bold": True, "color": GREEN, "nobullet": True}),
                 "So the negative result is informative, not an artifact of overly strict tests.",
                 ("", {"nobullet": True}),
                 ("4.  Implication for LLM-as-detector:", {"size": 16, "bold": True, "color": NAVY, "nobullet": True}),
                 "Same-dataset accuracy reports are uninformative.  Cross-dataset evaluation is necessary;  causal validation is desirable.",
                 ], font_size=14)

    add_textbox(s, 0.5, 6.5, 12.3, 0.7,
                "Code, data, paper draft:  github.com/yihuaihong/LLM_Rep_NetSec     ·     Thank you — questions?",
                font_size=14, bold=True, color=NAVY, align=PP_ALIGN.CENTER)
    add_footer(s, "16 / 16")

    # Save
    out_path = BASE / "LLM_RepSpace_NetSec_Presentation.pptx"
    prs.save(str(out_path))
    print(f"Saved: {out_path}")
    print(f"Slides: {len(prs.slides)}")
    print(f"Figures used:")
    for fp in sorted(DECK_FIGS.glob("*.png")):
        print(f"  {fp.name}")


if __name__ == "__main__":
    build_deck_figures()
    build_deck()
